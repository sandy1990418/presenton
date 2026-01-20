import os
from enum import Enum
from typing import Dict, List, Optional, Tuple, Union
from lxml import etree
from services.html_to_text_runs_service import (
    parse_html_text_to_text_runs as parse_inline_html_to_runs,
)

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.text.text import _Paragraph, TextFrame, Font, _Run
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.text import PP_ALIGN
from lxml.etree import fromstring, tostring
from PIL import Image
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

from pptx.util import Pt
from pptx.dml.color import RGBColor

from models.pptx_models import (
    PptxAutoShapeBoxModel,
    PptxBoxShapeEnum,
    PptxConnectorModel,
    PptxFillModel,
    PptxFontModel,
    PptxObjectFitEnum,
    PptxObjectFitModel,
    PptxParagraphModel,
    PptxPictureBoxModel,
    PptxPictureModel,
    PptxPositionModel,
    PptxPresentationModel,
    PptxShadowModel,
    PptxSlideModel,
    PptxSpacingModel,
    PptxStrokeModel,
    PptxStructureModel,
    PptxTableModel,
    PptxTextBoxModel,
    PptxTextRunModel,
)
from utils.download_helpers import download_files
from utils.image_utils import (
    clip_image,
    create_circle_image,
    fit_image,
    invert_image,
    round_image_corners,
    set_image_opacity,
)
import uuid

BLANK_SLIDE_LAYOUT = 6
TITLE_PLACEHOLDER_IDX = 0
CONTENT_PLACEHOLDER_IDX = 1

# 定義哪些 layout 需要分欄
MULTI_COLUMN_LAYOUTS: Dict[Union[str, int], int] = {
    "template_3": 2,
    "template_4": 3,
    "template_5": 4,
    3: 2,
    4: 3,
    5: 4,
}


# ==================== Image Layout ====================

class ImageLayoutEnum(Enum):
    # 右側
    RIGHT_HALF = "right_half"
    RIGHT_THIRD = "right_third"

    # 下方
    BOTTOM_HALF = "bottom_half"
    BOTTOM_TWO_THIRDS = "bottom_2_3"
    BOTTOM_THIRD = "bottom_third"

    # 全幅
    FULL = "full"


class ImageGridCalculator:
    """計算圖片 grid 位置"""

    def __init__(self, slide_width: float, slide_height: float, padding: float = 40):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.padding = padding

    def get_positions(
        self,
        image_count: int,
        layout: ImageLayoutEnum,
        content_top: float = 120,
    ) -> List[PptxPositionModel]:
        """根據 layout 計算每張圖片的位置"""

        if layout == ImageLayoutEnum.RIGHT_HALF:
            return self._vertical_grid(
                image_count,
                area_left=self.slide_width / 2 + self.padding / 2,
                area_top=self.padding,
                area_width=self.slide_width / 2 - self.padding * 1.5,
                area_height=self.slide_height - self.padding * 2,
            )

        elif layout == ImageLayoutEnum.RIGHT_THIRD:
            return self._vertical_grid(
                image_count,
                area_left=self.slide_width * 2 / 3 + self.padding / 2,
                area_top=self.padding,
                area_width=self.slide_width / 3 - self.padding * 1.5,
                area_height=self.slide_height - self.padding * 2,
            )

        elif layout == ImageLayoutEnum.BOTTOM_HALF:
            return self._horizontal_grid(
                image_count,
                area_left=self.padding,
                area_top=self.slide_height / 2 + self.padding / 2,
                area_width=self.slide_width - self.padding * 2,
                area_height=self.slide_height / 2 - self.padding * 1.5,
            )

        elif layout == ImageLayoutEnum.BOTTOM_TWO_THIRDS:
            text_area_height = (self.slide_height - content_top) / 3
            return self._horizontal_grid(
                image_count,
                area_left=self.padding,
                area_top=content_top + text_area_height + self.padding,
                area_width=self.slide_width - self.padding * 2,
                area_height=self.slide_height
                - (content_top + text_area_height)
                - self.padding * 2,
            )

        elif layout == ImageLayoutEnum.BOTTOM_THIRD:
            return self._horizontal_grid(
                image_count,
                area_left=self.padding,
                area_top=self.slide_height * 2 / 3 + self.padding / 2,
                area_width=self.slide_width - self.padding * 2,
                area_height=self.slide_height / 3 - self.padding * 1.5,
            )

        else:  # FULL
            return self._full_grid(image_count, content_top)

    def _vertical_grid(
        self,
        count: int,
        area_left: float,
        area_top: float,
        area_width: float,
        area_height: float,
    ) -> List[PptxPositionModel]:
        """垂直排列 (適合右側 layout)"""
        if count == 0:
            return []

        gap = self.padding / 2
        img_height = (area_height - gap * (count - 1)) / count

        return [
            PptxPositionModel(
                left=int(area_left),
                top=int(area_top + i * (img_height + gap)),
                width=int(area_width),
                height=int(img_height),
            )
            for i in range(count)
        ]

    def _horizontal_grid(
        self,
        count: int,
        area_left: float,
        area_top: float,
        area_width: float,
        area_height: float,
    ) -> List[PptxPositionModel]:
        """水平排列 (適合下方 layout)"""
        if count == 0:
            return []

        gap = self.padding / 2
        img_width = (area_width - gap * (count - 1)) / count

        return [
            PptxPositionModel(
                left=int(area_left + i * (img_width + gap)),
                top=int(area_top),
                width=int(img_width),
                height=int(area_height),
            )
            for i in range(count)
        ]

    def _full_grid(self, count: int, content_top: float) -> List[PptxPositionModel]:
        """全幅 grid 排列"""
        if count == 0:
            return []

        cols, rows = self._get_grid_size(count)

        area_top = content_top + self.padding
        area_height = self.slide_height - area_top - self.padding
        area_width = self.slide_width - self.padding * 2

        gap = self.padding / 2
        cell_width = (area_width - gap * (cols - 1)) / cols
        cell_height = (area_height - gap * (rows - 1)) / rows

        positions = []
        for i in range(count):
            row = i // cols
            col = i % cols
            positions.append(
                PptxPositionModel(
                    left=int(self.padding + col * (cell_width + gap)),
                    top=int(area_top + row * (cell_height + gap)),
                    width=int(cell_width),
                    height=int(cell_height),
                )
            )

        return positions

    def _get_grid_size(self, count: int) -> Tuple[int, int]:
        """決定 grid 的 cols x rows"""
        grid_map = {
            1: (1, 1),
            2: (2, 1),
            3: (3, 1),
            4: (2, 2),
            5: (3, 2),
            6: (3, 2),
            7: (4, 2),
            8: (4, 2),
            9: (3, 3),
        }
        return grid_map.get(count, (4, (count + 3) // 4))


# ==================== Main Creator ====================

class PptxPresentationCreator:
    """
    Create a PowerPoint (.pptx) presentation from a structured presentation model.
    """

    def __init__(
        self,
        ppt_model: PptxPresentationModel,
        temp_dir: str,
        template_path: Optional[str] = None,
    ):
        self._temp_dir = temp_dir
        self._ppt_model = ppt_model
        self._slide_models = ppt_model.slides

        # 存每個 slide 的 image layout 資訊
        self._slide_image_layouts: Dict[int, ImageLayoutEnum] = {}

        if template_path and os.path.exists(template_path):
            self._ppt = Presentation(template_path)
        else:
            self._ppt = Presentation()
            self._ppt.slide_width = Pt(1280)
            self._ppt.slide_height = Pt(720)

    # ==================== Class Method: from_simple_json ====================

    @classmethod
    def from_simple_json(
        cls,
        slides_data: List[dict],
        temp_dir: str,
        template_path: str,
        default_layout_index: int = 1,
    ) -> "PptxPresentationCreator":
        """
        從簡化 JSON 創建 PptxPresentationCreator

        JSON 格式:
        {
            "mainTitle": "標題",
            "layout_index": 1 或 "template_3",
            "bulletPoints": [{"text": "項目", "subPoints": ["子項目"]}],
            "image": {"__image_url__": "https://..."} 或 [{"__image_url__": "..."}],
            "imageLayout": "right_half" | "bottom_half" | ...,
            "__speaker_note__": "講者備註"
        }
        """
        template_ppt = Presentation(template_path)
        slide_width = template_ppt.slide_width.pt
        slide_height = template_ppt.slide_height.pt

        grid_calc = ImageGridCalculator(slide_width, slide_height)
        slides = []
        slide_image_layouts: Dict[int, ImageLayoutEnum] = {}

        for slide_index, slide_data in enumerate(slides_data):
            shapes = []

            # 取得 layout_index 和 column 數
            layout_key = slide_data.get("layout_index", default_layout_index)
            num_columns = MULTI_COLUMN_LAYOUTS.get(layout_key, 1)

            if isinstance(layout_key, int):
                layout_index = layout_key
            else:
                layout_index = default_layout_index

            # 收集圖片 - image 可能是 dict 或 list
            image_urls = []
            image_data = slide_data.get("image")

            if image_data:
                if isinstance(image_data, list):
                    image_urls = [
                        img.get("__image_url__")
                        for img in image_data
                        if img.get("__image_url__")
                    ]
                elif isinstance(image_data, dict):
                    if url := image_data.get("__image_url__"):
                        image_urls = [url]

            # 決定 image layout
            image_layout = None
            if image_urls:
                layout_str = slide_data.get("imageLayout")
                if layout_str:
                    image_layout = ImageLayoutEnum(layout_str)
                else:
                    image_layout = ImageLayoutEnum.RIGHT_HALF

                slide_image_layouts[slide_index] = image_layout

            # Title
            if main_title := slide_data.get("mainTitle"):
                shapes.append(
                    PptxTextBoxModel(
                        position=PptxPositionModel(left=0, top=0, width=0, height=0),
                        structure=PptxStructureModel(
                            level=0, isList=False, placeholder_idx=0
                        ),
                        paragraphs=[PptxParagraphModel(text=main_title)],
                    )
                )

            # Bullet points
            bullets = slide_data.get("bulletPoints", [])

            if bullets:
                if num_columns > 1:
                    # 多欄：平均分配 bullets
                    shapes.extend(
                        cls._distribute_bullets_to_columns(bullets, num_columns)
                    )
                else:
                    # 單欄：用 placeholder
                    for bp in bullets:
                        shapes.append(
                            PptxTextBoxModel(
                                position=PptxPositionModel(
                                    left=0, top=0, width=0, height=0
                                ),
                                structure=PptxStructureModel(
                                    level=0, isList=True, placeholder_idx=1
                                ),
                                paragraphs=[
                                    PptxParagraphModel(text=bp.get("text", ""))
                                ],
                            )
                        )
                        for sub in bp.get("subPoints", []):
                            shapes.append(
                                PptxTextBoxModel(
                                    position=PptxPositionModel(
                                        left=0, top=0, width=0, height=0
                                    ),
                                    structure=PptxStructureModel(
                                        level=1, isList=True, placeholder_idx=1
                                    ),
                                    paragraphs=[PptxParagraphModel(text=sub)],
                                )
                            )

            # Images
            if image_urls and image_layout:
                positions = grid_calc.get_positions(len(image_urls), image_layout)

                for img_url, pos in zip(image_urls, positions):
                    shapes.append(
                        PptxPictureBoxModel(
                            position=pos,
                            picture=PptxPictureModel(
                                path=img_url,
                                is_network=img_url.startswith("http"),
                            ),
                            object_fit=PptxObjectFitModel(fit=PptxObjectFitEnum.CONTAIN),
                            clip=False,
                        )
                    )

            slides.append(
                PptxSlideModel(
                    shapes=shapes,
                    layout_index=layout_index,
                    note=slide_data.get("__speaker_note__"),
                )
            )

        ppt_model = PptxPresentationModel(slides=slides)
        creator = cls(ppt_model, temp_dir, template_path)
        creator._slide_image_layouts = slide_image_layouts

        return creator

    @classmethod
    def _distribute_bullets_to_columns(
        cls,
        bullets: List[dict],
        num_columns: int,
    ) -> List[PptxTextBoxModel]:
        """將 bullets 平均分配到多欄"""
        shapes = []

        bullets_per_column = len(bullets) // num_columns
        remainder = len(bullets) % num_columns

        bullet_idx = 0
        for col in range(num_columns):
            count = bullets_per_column + (1 if col < remainder else 0)
            column_bullets = bullets[bullet_idx : bullet_idx + count]
            bullet_idx += count

            for bp in column_bullets:
                # 主項目
                shapes.append(
                    PptxTextBoxModel(
                        position=PptxPositionModel(left=0, top=0, width=0, height=0),
                        structure=PptxStructureModel(
                            level=0,
                            isList=True,
                            placeholder_idx=1,
                            column=col,
                            num_columns=num_columns,
                        ),
                        paragraphs=[PptxParagraphModel(text=bp.get("text", ""))],
                    )
                )
                # 子項目
                for sub in bp.get("subPoints", []):
                    shapes.append(
                        PptxTextBoxModel(
                            position=PptxPositionModel(left=0, top=0, width=0, height=0),
                            structure=PptxStructureModel(
                                level=1,
                                isList=True,
                                placeholder_idx=1,
                                column=col,
                                num_columns=num_columns,
                            ),
                            paragraphs=[PptxParagraphModel(text=sub)],
                        )
                    )

        return shapes

    # ==================== Main Flow ====================

    async def create_ppt(self):
        """Build the presentation content in memory."""
        await self.fetch_network_assets()

        for index, slide_model in enumerate(self._slide_models):
            if self._ppt_model.shapes:
                slide_model.shapes.extend(self._ppt_model.shapes)

            self._add_slide(slide_model, index)

    def _add_slide(self, slide_model: PptxSlideModel, slide_index: int):
        """Create and populate a single slide."""
        layout_index = slide_model.layout_index
        slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[layout_index])

        if slide_model.background:
            self._apply_fill(slide.background, slide_model.background)

        if slide_model.note:
            slide.notes_slide.notes_text_frame.text = slide_model.note

        # 取得這個 slide 的 image layout
        image_layout = self._slide_image_layouts.get(slide_index)

        # 分類 shapes
        title_textboxes: List[PptxTextBoxModel] = []
        content_textboxes: List[PptxTextBoxModel] = []  # 單欄
        column_textboxes: Dict[int, List[PptxTextBoxModel]] = {}  # 多欄
        free_shapes = []

        for shape in slide_model.shapes:
            if isinstance(shape, PptxTextBoxModel) and shape.structure:
                if shape.structure.placeholder_idx == TITLE_PLACEHOLDER_IDX:
                    title_textboxes.append(shape)
                elif shape.structure.column is not None:
                    col = shape.structure.column
                    if col not in column_textboxes:
                        column_textboxes[col] = []
                    column_textboxes[col].append(shape)
                else:
                    content_textboxes.append(shape)
            else:
                free_shapes.append(shape)

        # 填入 title placeholder
        if title_textboxes:
            title_ph = self._get_placeholder(slide, TITLE_PLACEHOLDER_IDX)
            if title_ph:
                self._fill_placeholder(title_ph.text_frame, title_textboxes)
            else:
                free_shapes.extend(title_textboxes)

        # 填入 content placeholder（單欄）
        if content_textboxes:
            content_ph = self._get_placeholder(slide, CONTENT_PLACEHOLDER_IDX)
            if content_ph:
                self._adjust_placeholder_for_layout(content_ph, image_layout)
                self._fill_placeholder(content_ph.text_frame, content_textboxes)
            else:
                free_shapes.extend(content_textboxes)

        # 處理多欄
        if column_textboxes:
            self._add_column_textboxes(slide, column_textboxes, image_layout)

        # 加入 free shapes
        for shape in free_shapes:
            self._add_shape(slide, shape)

    # ==================== Placeholder ====================

    def _get_placeholder(self, slide: Slide, idx: int):
        """Retrieve a placeholder by index from a slide."""
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == idx:
                return shape
        return None

    def _fill_placeholder(
        self, text_frame: TextFrame, textboxes: List[PptxTextBoxModel]
    ):
        """Populate a placeholder text frame while preserving template bullet styles."""
        first_para = True

        for textbox in textboxes:
            level = textbox.structure.level
            is_list = textbox.structure.isList

            for para_model in textbox.paragraphs:
                if first_para:
                    para = text_frame.paragraphs[0]
                    first_para = False
                else:
                    para = text_frame.add_paragraph()

                para.level = level
                para.alignment = PP_ALIGN.LEFT

                if not is_list:
                    self._disable_bullet(para)

                self._populate_paragraph(para, para_model)

    def _adjust_placeholder_for_layout(
        self, placeholder, image_layout: Optional[ImageLayoutEnum]
    ):
        """根據 image layout 調整 placeholder"""
        if not image_layout:
            return

        slide_width = self._ppt.slide_width.pt
        slide_height = self._ppt.slide_height.pt
        padding = 40

        # 先保存原始值
        original_width = placeholder.width
        original_height = placeholder.height
        ph_top = placeholder.top.pt if placeholder.top else 120

        if image_layout == ImageLayoutEnum.RIGHT_HALF:
            placeholder.width = Pt(slide_width / 2 - padding * 2)

        elif image_layout == ImageLayoutEnum.RIGHT_THIRD:
            placeholder.width = Pt(slide_width * 2 / 3 - padding * 2)

        elif image_layout in [
            ImageLayoutEnum.BOTTOM_HALF,
            ImageLayoutEnum.BOTTOM_TWO_THIRDS,
            ImageLayoutEnum.BOTTOM_THIRD,
        ]:
            # 計算圖片區域的 top 位置
            if image_layout == ImageLayoutEnum.BOTTOM_HALF:
                image_area_top = slide_height / 2
            elif image_layout == ImageLayoutEnum.BOTTOM_TWO_THIRDS:
                image_area_top = slide_height / 3
            else:  # BOTTOM_THIRD
                image_area_top = slide_height * 2 / 3

            # 保持原始 width，調整 height
            new_height = image_area_top - ph_top - padding
            placeholder.height = Pt(max(new_height, 100))
            placeholder.width = original_width  # 恢復原始 width

    def _disable_bullet(self, paragraph: _Paragraph):
        """Disable bullet formatting for a paragraph."""
        pPr = paragraph._p.get_or_add_pPr()

        for tag in ["a:buFont", "a:buChar", "a:buAutoNum"]:
            elem = pPr.find(qn(tag))
            if elem is not None:
                pPr.remove(elem)

        if pPr.find(qn("a:buNone")) is None:
            pPr.append(OxmlElement("a:buNone"))

    def _enable_bullet(self, paragraph: _Paragraph, level: int = 0):
        """為段落啟用 bullet"""
        pPr = paragraph._p.get_or_add_pPr()

        # 移除 buNone
        buNone = pPr.find(qn("a:buNone"))
        if buNone is not None:
            pPr.remove(buNone)

        # 設定 bullet 字符（根據 level）
        bullet_chars = ["•", "–", "◦", "▪"]
        bullet_char = bullet_chars[min(level, len(bullet_chars) - 1)]

        # 移除舊的 buChar
        old_buChar = pPr.find(qn("a:buChar"))
        if old_buChar is not None:
            pPr.remove(old_buChar)

        # 加入新的 buChar
        buChar = OxmlElement("a:buChar")
        buChar.set("char", bullet_char)
        pPr.append(buChar)

    # ==================== Multi-Column ====================

    def _add_column_textboxes(
        self,
        slide: Slide,
        column_textboxes: Dict[int, List[PptxTextBoxModel]],
        image_layout: Optional[ImageLayoutEnum],
    ):
        """建立多欄 textbox，保留 bullet 樣式"""
        num_columns = max(column_textboxes.keys()) + 1 if column_textboxes else 0
        if num_columns == 0:
            return

        slide_width = self._ppt.slide_width.pt
        slide_height = self._ppt.slide_height.pt
        padding = 40
        content_top = 120

        # 根據 image_layout 調整可用寬度
        if image_layout == ImageLayoutEnum.RIGHT_HALF:
            available_width = slide_width / 2 - padding * 2
        elif image_layout == ImageLayoutEnum.RIGHT_THIRD:
            available_width = slide_width * 2 / 3 - padding * 2
        else:
            available_width = slide_width - padding * 2

        # 根據 image_layout 調整高度
        if image_layout == ImageLayoutEnum.BOTTOM_HALF:
            content_height = slide_height / 2 - content_top - padding
        elif image_layout == ImageLayoutEnum.BOTTOM_TWO_THIRDS:
            content_height = slide_height / 3 - padding
        elif image_layout == ImageLayoutEnum.BOTTOM_THIRD:
            content_height = slide_height * 2 / 3 - content_top - padding
        else:
            content_height = slide_height - content_top - padding * 2

        gap = 20
        column_width = (available_width - gap * (num_columns - 1)) / num_columns

        for col_idx in sorted(column_textboxes.keys()):
            col_items = column_textboxes[col_idx]
            col_left = padding + col_idx * (column_width + gap)

            # 創建 textbox
            shape = slide.shapes.add_textbox(
                Pt(col_left), Pt(content_top), Pt(column_width), Pt(content_height)
            )
            tf = shape.text_frame
            tf.word_wrap = True

            # 填入內容
            first_para = True
            for textbox in col_items:
                level = textbox.structure.level
                is_list = textbox.structure.isList

                for para_model in textbox.paragraphs:
                    if first_para:
                        para = tf.paragraphs[0]
                        first_para = False
                    else:
                        para = tf.add_paragraph()

                    para.level = level

                    if is_list:
                        self._enable_bullet(para, level)
                    else:
                        self._disable_bullet(para)

                    self._populate_paragraph(para, para_model)

    # ==================== Free-position Shapes ====================

    def _add_shape(self, slide: Slide, shape):
        """Dispatch shape creation based on its model type."""
        if isinstance(shape, PptxPictureBoxModel):
            self._add_picture(slide, shape)
        elif isinstance(shape, PptxAutoShapeBoxModel):
            self._add_autoshape(slide, shape)
        elif isinstance(shape, PptxTextBoxModel):
            self._add_textbox(slide, shape)
        elif isinstance(shape, PptxConnectorModel):
            self._add_connector(slide, shape)
        elif isinstance(shape, PptxTableModel):
            self._add_table(slide, shape)

    def _add_textbox(self, slide: Slide, model: PptxTextBoxModel):
        """Add a textbox shape to a slide."""
        shape = slide.shapes.add_textbox(*model.position.to_pt_list())
        shape.width += Pt(2)

        tf = shape.text_frame
        tf.word_wrap = model.text_wrap

        self._apply_fill(shape, model.fill)
        self._apply_margin(tf, model.margin)
        self._add_paragraphs(tf, model.paragraphs)

    def _add_autoshape(self, slide: Slide, model: PptxAutoShapeBoxModel):
        """Add an auto shape to a slide."""
        position = self._get_margined_position(model.position, model.margin)
        shape = slide.shapes.add_shape(model.type, *position.to_pt_list())

        tf = shape.text_frame
        tf.word_wrap = model.text_wrap

        self._apply_fill(shape, model.fill)
        self._apply_margin(tf, model.margin)
        self._apply_stroke(shape, model.stroke)
        self._apply_shadow(shape, model.shadow)
        self._apply_border_radius(shape, model.border_radius)

        if model.paragraphs:
            self._add_paragraphs(tf, model.paragraphs)

    def _add_picture(self, slide: Slide, model: PptxPictureBoxModel):
        """Add a picture shape to a slide."""
        image_path = model.picture.path
        final_position = self._get_margined_position(model.position, model.margin)

        needs_processing = (
            model.clip
            or model.border_radius
            or model.invert
            or model.opacity
            or (model.object_fit and model.object_fit.fit)
            or model.shape
        )

        if needs_processing:
            try:
                image = Image.open(image_path).convert("RGBA")
            except Exception:
                return

            if model.object_fit and model.object_fit.fit:
                image = fit_image(
                    image,
                    final_position.width,
                    final_position.height,
                    model.object_fit,
                )
            elif model.clip:
                image = clip_image(
                    image,
                    final_position.width,
                    final_position.height,
                )

            if model.border_radius:
                image = round_image_corners(image, model.border_radius)
            if model.shape == PptxBoxShapeEnum.CIRCLE:
                image = create_circle_image(image)
            if model.invert:
                image = invert_image(image)
            if model.opacity:
                image = set_image_opacity(image, model.opacity)

            image_path = os.path.join(self._temp_dir, f"{uuid.uuid4()}.png")
            image.save(image_path)

        slide.shapes.add_picture(image_path, *final_position.to_pt_list())

    def _add_connector(self, slide: Slide, model: PptxConnectorModel):
        """Add a connector (line) shape to a slide."""
        if model.thickness == 0:
            return
        shape = slide.shapes.add_connector(model.type, *model.position.to_pt_xyxy())
        shape.line.width = Pt(model.thickness)
        shape.line.color.rgb = RGBColor.from_string(model.color)
        self._set_fill_opacity(shape, model.opacity)

    def _add_table(self, slide: Slide, model: PptxTableModel):
        """Add a table shape to a slide."""
        if not model.rows:
            return

        num_rows = len(model.rows)
        num_cols = len(model.rows[0]) if model.rows else 0

        if num_rows == 0 or num_cols == 0:
            return

        table_shape = slide.shapes.add_table(
            num_rows,
            num_cols,
            Pt(model.position.left),
            Pt(model.position.top),
            Pt(model.position.width),
            Pt(model.position.height),
        )
        table = table_shape.table

        default_font = model.font or PptxFontModel(
            name="Inter", size=12, color="000000"
        )
        header_font = model.header_font or PptxFontModel(
            name="Inter", size=12, color="FFFFFF", font_weight=700
        )

        for row_idx, row in enumerate(model.rows):
            is_header = model.header_row and row_idx == 0

            for col_idx, cell_model in enumerate(row):
                if col_idx >= num_cols:
                    break

                cell = table.cell(row_idx, col_idx)
                cell.text = cell_model.text

                if is_header and model.header_fill:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(
                        model.header_fill.color
                    )
                elif not is_header and model.cell_fill:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(
                        model.cell_fill.color
                    )

                font_model = header_font if is_header else default_font
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = font_model.name
                    paragraph.font.size = Pt(font_model.size)
                    paragraph.font.color.rgb = RGBColor.from_string(font_model.color)
                    paragraph.font.bold = (
                        font_model.font_weight >= 600 if font_model.font_weight else False
                    )

    # ==================== Paragraph / TextRun ====================

    def _add_paragraphs(self, text_frame: TextFrame, models: List[PptxParagraphModel]):
        """Add paragraphs to a text frame."""
        for i, model in enumerate(models):
            para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            self._populate_paragraph(para, model)

    def _populate_paragraph(self, para: _Paragraph, model: PptxParagraphModel):
        """Populate a paragraph with text runs and styles."""
        if model.spacing:
            para.space_before = Pt(model.spacing.top)
            para.space_after = Pt(model.spacing.bottom)

        if model.line_height:
            para.line_spacing = model.line_height

        if model.alignment:
            para.alignment = model.alignment

        if model.font:
            self._apply_font(para.font, model.font)

        text_runs = []
        if model.text:
            text_runs = parse_inline_html_to_runs(model.text, model.font)
        elif model.text_runs:
            text_runs = model.text_runs

        for run_model in text_runs:
            run = para.add_run()
            run.text = run_model.text
            if run_model.font:
                self._apply_font(run.font, run_model.font)

    # ==================== Style Helpers ====================

    def _apply_font(self, font: Font, model: PptxFontModel):
        """Apply font styling to a text element."""
        font.name = model.name
        font.size = Pt(model.size)
        font.color.rgb = RGBColor.from_string(model.color)
        font.italic = model.italic
        font.bold = model.font_weight >= 600

        if model.underline is not None:
            font.underline = bool(model.underline)
        if model.strike is not None:
            rPr = font._element
            rPr.set("strike", "sngStrike" if model.strike else "noStrike")

    def _apply_fill(self, shape: Shape, fill: Optional[PptxFillModel]):
        """Apply fill color and opacity to a shape."""
        if not fill:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.color)
            self._set_fill_opacity(shape.fill, fill.opacity)

    def _apply_stroke(self, shape: Shape, stroke: Optional[PptxStrokeModel]):
        """Apply stroke (outline) styling to a shape."""
        if not stroke or stroke.thickness == 0:
            shape.line.fill.background()
        else:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = RGBColor.from_string(stroke.color)
            shape.line.width = Pt(stroke.thickness)
            self._set_fill_opacity(shape.line.fill, stroke.opacity)

    def _apply_shadow(self, shape: Shape, shadow: Optional[PptxShadowModel]):
        """Apply shadow effect to a shape."""
        sp_pr = shape._element.xpath("p:spPr")[0]
        nsmap = sp_pr.nsmap

        effect_list = sp_pr.find("a:effectLst", namespaces=nsmap)
        if effect_list is not None:
            for tag in ["a:outerShdw", "a:innerShdw", "a:prstShdw"]:
                old = effect_list.find(tag, namespaces=nsmap)
                if old is not None:
                    effect_list.remove(old)
        else:
            effect_list = etree.SubElement(
                sp_pr,
                f"{{{nsmap['a']}}}effectLst",
                nsmap=nsmap,
            )

        if shadow is None:
            outer = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                {"blurRad": "0", "dist": "0", "dir": "0"},
                nsmap=nsmap,
            )
            color = etree.SubElement(
                outer,
                f"{{{nsmap['a']}}}srgbClr",
                {"val": "000000"},
                nsmap=nsmap,
            )
            etree.SubElement(
                color,
                f"{{{nsmap['a']}}}alpha",
                {"val": "0"},
                nsmap=nsmap,
            )
        else:
            angle = int(round((shadow.angle % 360) * 60000)) if shadow.angle else 0
            outer = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                {
                    "blurRad": f"{Pt(shadow.radius)}",
                    "dir": f"{angle}",
                    "dist": f"{Pt(shadow.offset)}",
                    "rotWithShape": "0",
                },
                nsmap=nsmap,
            )
            color = etree.SubElement(
                outer,
                f"{{{nsmap['a']}}}srgbClr",
                {"val": shadow.color},
                nsmap=nsmap,
            )
            etree.SubElement(
                color,
                f"{{{nsmap['a']}}}alpha",
                {"val": f"{int(shadow.opacity * 100000)}"},
                nsmap=nsmap,
            )

    def _apply_border_radius(self, shape: Shape, radius: Optional[int]):
        """Apply rounded corners to a shape if supported."""
        if not radius:
            return
        try:
            shape.adjustments[0] = Pt(radius) / min(shape.width, shape.height)
        except Exception:
            pass

    def _apply_margin(self, tf: TextFrame, margin: Optional[PptxSpacingModel]):
        """Apply text margins to a text frame."""
        tf.margin_left = Pt(margin.left if margin else 0)
        tf.margin_right = Pt(margin.right if margin else 0)
        tf.margin_top = Pt(margin.top if margin else 0)
        tf.margin_bottom = Pt(margin.bottom if margin else 0)

    def _set_fill_opacity(self, fill, opacity):
        """Set opacity on a fill element."""
        if opacity is None or opacity >= 1.0:
            return
        try:
            sF = fill._xPr.solidFill.get_or_change_to_srgbClr()
            elem = OxmlElement("a:alpha")
            elem.set("val", str(int(opacity * 100000)))
            sF.append(elem)
        except Exception:
            pass

    def _get_margined_position(
        self,
        pos: PptxPositionModel,
        margin: Optional[PptxSpacingModel],
    ) -> PptxPositionModel:
        """Compute a position adjusted by margins."""
        if not margin:
            return pos
        return PptxPositionModel(
            left=pos.left + margin.left,
            top=pos.top + margin.top,
            width=max(pos.width - margin.left - margin.right, 0),
            height=max(pos.height - margin.top - margin.bottom, 0),
        )

    # ==================== Network Assets ====================

    async def fetch_network_assets(self):
        """Download all remote image assets referenced by the presentation model."""
        image_urls = []
        models_with_network_asset: List[PptxPictureBoxModel] = []

        all_shapes = list(self._ppt_model.shapes or [])
        for slide in self._slide_models:
            all_shapes.extend(slide.shapes)

        for shape in all_shapes:
            if isinstance(shape, PptxPictureBoxModel):
                path = shape.picture.path
                if path.startswith("http"):
                    if "app_data" in path:
                        shape.picture.path = os.path.join(
                            "/app_data",
                            path.split("app_data/")[1],
                        )
                        shape.picture.is_network = False
                    else:
                        image_urls.append(path)
                        models_with_network_asset.append(shape)

        if image_urls:
            paths = await download_files(image_urls, self._temp_dir)
            for shape, path in zip(models_with_network_asset, paths):
                if path:
                    shape.picture.path = path
                    shape.picture.is_network = False

    # ==================== Save ====================

    def save(self, path: str):
        """Save the generated presentation to disk."""
        self._ppt.save(path)